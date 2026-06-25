"""Helpers de design pour le PPT collectif Torpier — palette charte graphique."""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# Palette Torpier
ORANGE   = RGBColor(0xE9, 0x55, 0x13)
ANTHRA   = RGBColor(0x3C, 0x3C, 0x3B)
GRIS     = RGBColor(0xD6, 0xCD, 0xC6)
CREME    = RGBColor(0xF8, 0xF0, 0xEA)
BLANC    = RGBColor(0xFF, 0xFF, 0xFF)
NOIR     = RGBColor(0x1D, 0x1D, 0x1B)
# Teintes sémantiques (SWOT, statuts)
VERT     = RGBColor(0x4E, 0x7A, 0x3F)
ROUGE    = RGBColor(0xB0, 0x3A, 0x2E)

SW = Inches(10)        # largeur slide
SH = Inches(5.625)     # hauteur slide

def blank_slide(prs):
    """Crée une slide vierge (fond blanc) et renvoie l'objet slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[10])  # BLANK
    # fond blanc explicite
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = BLANC
    bg.line.fill.background()
    bg.shadow.inherit = False
    _send_back(slide, bg)
    return slide

def _send_back(slide, shape):
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)

def rect(slide, x, y, w, h, color, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(1)
    s.shadow.inherit = False
    return s

def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    return tb, tf

def set_para(p, text, size=12, color=ANTHRA, bold=False, align=PP_ALIGN.LEFT,
             space_after=4, space_before=0, font="Arial"):
    p.text = text
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    for r in p.runs:
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = font
    return p

def add_para(tf, text, **kw):
    p = tf.add_paragraph()
    return set_para(p, text, **kw)

def header(slide, title, subtitle=None):
    """Bandeau de titre standard : accent orange + titre anthracite + hairline."""
    rect(slide, Inches(0.5), Inches(0.42), Inches(0.10), Inches(0.55), ORANGE)
    tb, tf = textbox(slide, Inches(0.72), Inches(0.38), Inches(8.8), Inches(0.7),
                     anchor=MSO_ANCHOR.MIDDLE)
    set_para(tf.paragraphs[0], title, size=24, color=ANTHRA, bold=True)
    if subtitle:
        add_para(tf, subtitle, size=12, color=ORANGE, bold=True, space_before=2)
    # hairline
    rect(slide, Inches(0.5), Inches(1.18), Inches(9.0), Pt(1.5), GRIS)

def footer(slide, n):
    tb, tf = textbox(slide, Inches(8.7), Inches(5.28), Inches(1.1), Inches(0.3))
    set_para(tf.paragraphs[0], f"Torpier · {n}", size=8, color=GRIS, align=PP_ALIGN.RIGHT)

def card(slide, x, y, w, h, fill=CREME):
    """Carte de fond arrondie pour regrouper du contenu."""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    try:
        s.adjustments[0] = 0.05
    except Exception:
        pass
    return s

def chip(slide, x, y, w, text, color=ORANGE, textcolor=BLANC, h=Inches(0.34)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    set_para(tf.paragraphs[0], text, size=12.5, color=textcolor, bold=True,
             align=PP_ALIGN.CENTER)
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return s

def notes(slide, text):
    """Renseigne les notes de l'orateur de la slide."""
    slide.notes_slide.notes_text_frame.text = text.strip()


def remove_example_slides(prs, keep_indexes):
    """Supprime toutes les slides sauf celles d'indice dans keep_indexes."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    for i, sld in enumerate(slides):
        if i not in keep_indexes:
            xml_slides.remove(sld)
