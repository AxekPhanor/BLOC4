# -*- coding: utf-8 -*-
"""Genere les 2 supports de soutenance BLOC4 (Torpier) : individuel + collectif."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- Charte ----
NAVY   = RGBColor(0x16, 0x2A, 0x4A)   # titres
TEAL   = RGBColor(0x00, 0x9B, 0x91)   # accent CESI
LIGHT  = RGBColor(0xF2, 0xF5, 0xF7)   # fond bloc
GREY   = RGBColor(0x55, 0x5B, 0x66)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xE8, 0x7A, 0x1E)
DARK   = RGBColor(0x22, 0x28, 0x33)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)   # 16:9


def new_prez():
    p = Presentation()
    p.slide_width = EMU_W
    p.slide_height = EMU_H
    return p


def _blank(prez):
    return prez.slides.add_slide(prez.slide_layouts[6])


def band(slide, color, x, y, w, h):
    shp = slide.shapes.add_shape(1, x, y, w, h)  # rectangle
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def set_run(r, text, size, color, bold=False, italic=False, font="Calibri"):
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font


# ---------- Slides types ----------

def title_slide(prez, kicker, title, subtitle, footer):
    s = _blank(prez)
    band(s, NAVY, 0, 0, EMU_W, EMU_H)
    band(s, TEAL, 0, Inches(4.55), EMU_W, Inches(0.12))
    # kicker
    _, tf = textbox(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(0.6))
    set_run(tf.paragraphs[0].add_run(), kicker, 18, TEAL, bold=True)
    # title
    _, tf = textbox(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.7))
    set_run(tf.paragraphs[0].add_run(), title, 40, WHITE, bold=True)
    # subtitle
    _, tf = textbox(s, Inches(0.9), Inches(4.8), Inches(11.5), Inches(1.2))
    for i, line in enumerate(subtitle):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        set_run(p.add_run(), line, 18, RGBColor(0xC9, 0xD4, 0xE0))
    # footer
    _, tf = textbox(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.5))
    set_run(tf.paragraphs[0].add_run(), footer, 12, RGBColor(0x9A, 0xA7, 0xB5))
    return s


def header(slide, theme, title, pts=None):
    band(slide, WHITE, 0, 0, EMU_W, EMU_H)
    band(slide, NAVY, 0, 0, EMU_W, Inches(1.15))
    band(slide, TEAL, 0, Inches(1.15), EMU_W, Inches(0.06))
    _, tf = textbox(slide, Inches(0.6), Inches(0.12), Inches(10.6), Inches(0.45))
    set_run(tf.paragraphs[0].add_run(), theme, 13, TEAL, bold=True)
    _, tf = textbox(slide, Inches(0.6), Inches(0.5), Inches(10.6), Inches(0.62))
    set_run(tf.paragraphs[0].add_run(), title, 26, WHITE, bold=True)
    if pts is not None:
        chip = band(slide, TEAL, Inches(11.7), Inches(0.32), Inches(1.25), Inches(0.55))
        chip.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        chip.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_run(chip.text_frame.paragraphs[0].add_run(), pts, 14, WHITE, bold=True)


def footer(slide, idx, total, label):
    _, tf = textbox(slide, Inches(0.6), Inches(7.05), Inches(9), Inches(0.4))
    set_run(tf.paragraphs[0].add_run(), label, 9, GREY)
    _, tf = textbox(slide, Inches(11.8), Inches(7.05), Inches(1.2), Inches(0.4))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    set_run(p.add_run(), f"{idx} / {total}", 9, GREY)


def bullets(slide, items, x=Inches(0.7), y=Inches(1.5), w=Inches(12), h=Inches(5.2),
            size=17, gap=8):
    _, tf = textbox(slide, x, y, w, h)
    first = True
    for it in items:
        lvl = it[0] if isinstance(it, tuple) else 0
        txt = it[1] if isinstance(it, tuple) else it
        bold = it[2] if isinstance(it, tuple) and len(it) > 2 else False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(gap)
        marker = "▸ " if lvl == 0 else "•  "
        r = p.add_run()
        col = NAVY if lvl == 0 else GREY
        set_run(r, marker + txt, size - lvl * 2, col, bold=(bold or lvl == 0))


def table(slide, rows, x, y, w, col_w=None, header_fill=NAVY, fsize=12,
          h=None):
    nrows = len(rows); ncols = len(rows[0])
    if h is None:
        h = Inches(0.42 * nrows)
    gfx = slide.shapes.add_table(nrows, ncols, x, y, w, h)
    tbl = gfx.table
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Emu(int(w * cw / total))
    for r in range(nrows):
        for c in range(ncols):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            para = tf.paragraphs[0]
            run = para.add_run()
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
                set_run(run, rows[r][c], fsize, WHITE, bold=True)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
                set_run(run, rows[r][c], fsize, DARK)
    return tbl


def section(prez, n, title, sub=""):
    s = _blank(prez)
    band(s, NAVY, 0, 0, EMU_W, EMU_H)
    band(s, TEAL, Inches(0.9), Inches(3.05), Inches(0.9), Inches(0.12))
    _, tf = textbox(s, Inches(0.9), Inches(2.0), Inches(2), Inches(1))
    set_run(tf.paragraphs[0].add_run(), n, 60, TEAL, bold=True)
    _, tf = textbox(s, Inches(0.9), Inches(3.4), Inches(11), Inches(1.2))
    set_run(tf.paragraphs[0].add_run(), title, 32, WHITE, bold=True)
    if sub:
        _, tf = textbox(s, Inches(0.9), Inches(4.5), Inches(11), Inches(1))
        set_run(tf.paragraphs[0].add_run(), sub, 16, RGBColor(0xC9, 0xD4, 0xE0), italic=True)
    return s


# =========================================================================
#  PRESENTATION COLLECTIVE
# =========================================================================

def build_collective(path):
    P = new_prez()
    FOOT = "BLOC4 · Manager les equipes et la transformation du SI · Projet collaboratif (cas Torpier)"
    total = 14
    i = [0]

    def foot(s):
        i[0] += 1
        footer(s, i[0], total, FOOT)

    # 1 Title
    s = title_slide(P, "PROJET EXPERIENTIEL COLLABORATIF",
                    "Manager les equipes et la transformation du SI",
                    ["Etude preliminaire d'evolution du SI — ebauche de schema directeur",
                     "Groupe DSI · Soutenance 20 min + 15 min d'echange"],
                    "Cas entreprise TORPIER — CESI")
    foot(s)

    # 2 Sommaire / organisation du groupe
    s = _blank(P); header(s, "Introduction", "Sommaire & organisation du groupe", "1 pt")
    bullets(s, [
        (0, "Objectif : etude preliminaire d'evolution du SI alignee sur la strategie Torpier", True),
        (1, "1. Analyse de la situation initiale (SWOT, processus, applications, infrastructure)"),
        (1, "2. 5 axes d'amelioration + comparaison de 3 referentiels de gouvernance"),
        (1, "3. Organisation des equipes, RACI et outils de pilotage"),
        (1, "4. Conduite du changement (pitch d'adhesion) + bilan du groupe"),
        (0, "Organisation du groupe : repartition des roles sur l'etude", True),
    ], y=Inches(1.45), h=Inches(3.2))
    table(s, [
        ["Membre", "Role sur l'etude", "Livrable principal"],
        ["Membre 1 (pilote)", "Coordination + gouvernance", "Comparaison referentiels"],
        ["Membre 2", "Cartographie applicative & technique", "Tableaux applis / infra"],
        ["Membre 3", "Processus metiers & SWOT", "Cartographie processus"],
        ["Membre 4", "Organisation & conduite du changement", "RACI + pitch"],
    ], Inches(0.7), Inches(4.7), Inches(12), col_w=[2, 3, 3], fsize=12)
    foot(s)

    # 3 SWOT
    s = _blank(P); header(s, "1. Modeliser les flux metiers et ressources techniques",
                          "SWOT de l'entreprise Torpier", "1 pt")
    quad = [
        ("FORCES", TEAL, ["Savoir-faire bois & eco-conception", "CA +30% (2023-25), expansion EU",
                          "Applications metiers internes (DSI)", "E-commerce PrestaShop + SAV bots"]),
        ("FAIBLESSES", ORANGE, ["Silos applicatifs (echanges CSV)", "Exchange en migration, dette technique",
                                "Pilotage projets SI 'en construction'", "Heterogeneite FR / Finlande"]),
        ("OPPORTUNITES", RGBColor(0x2E,0x7D,0x46), ["Industrie 4.0 / IoT en production", "Cloud hybride & M365",
                                "Clients prives / promoteurs", "Demarche ITIL de services"]),
        ("MENACES", RGBColor(0xB0,0x3A,0x2E), ["Cyber-risques / RGPD", "Dependance fournisseurs bois",
                                "Resilience interconnexion sites", "Barriere linguistique/culture"]),
    ]
    xs = [Inches(0.7), Inches(6.85)]; ys = [Inches(1.45), Inches(4.25)]
    k = 0
    for (lbl, col, lst) in quad:
        bx = xs[k % 2]; by = ys[k // 2]
        box = band(s, LIGHT, bx, by, Inches(5.78), Inches(2.6))
        bar = band(s, col, bx, by, Inches(5.78), Inches(0.45))
        bar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        bar.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        set_run(bar.text_frame.paragraphs[0].add_run(), lbl, 14, WHITE, bold=True)
        _, tf = textbox(s, bx + Inches(0.15), by + Inches(0.5), Inches(5.5), Inches(2.0))
        for j, t in enumerate(lst):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_after = Pt(3)
            set_run(p.add_run(), "• " + t, 12.5, DARK)
        k += 1
    foot(s)

    # 4 Processus metiers
    s = _blank(P); header(s, "1. Modeliser", "Processus metiers supportes par le SI", "2 pts")
    table(s, [
        ["Processus", "Taches principales", "Acteurs cles"],
        ["Gestion commerciale & relation client", "Ventes, commandes, livraison, reclamations", "Commerciaux, SAV, Dynamics"],
        ["Production", "Approvisionnement, fabrication, conception", "Poles Arras/Finlande, conception"],
        ["Financier & administratif", "Comptabilite, facturation, paie, RH", "Administration, SAGE"],
        ["Gouvernance", "Strategie, KPI, arbitrage, risques", "Direction generale"],
        ["Supports IT", "Support, maintenance, securite, projets SI", "DSI"],
    ], Inches(0.7), Inches(1.55), Inches(12), col_w=[3, 4, 3], fsize=12.5)
    _, tf = textbox(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.8))
    set_run(tf.paragraphs[0].add_run(),
            "Representation conseillee : cartographie visuelle (chaine de valeur / BPMN) reliant processus, acteurs et flux d'information.",
            13, GREY, italic=True)
    foot(s)

    # 5 Applications
    s = _blank(P); header(s, "1. Modeliser", "Strate applicative — applis ↔ fonctions metiers", "2 pts")
    table(s, [
        ["Application", "Editeur / type", "Fonctions metiers servies"],
        ["Microsoft Dynamics", "Microsoft / Azure", "CRM, commandes, relation client"],
        ["GesProd (interne)", "Torpier / serveur", "Approvisionnement, stocks"],
        ["SAGE Compta / Paie", "SAGE / serveur", "Comptabilite, facturation, RH/paie"],
        ["Qualeval (interne)", "Torpier / serveur", "Qualite production, conformite"],
        ["PrestaShop (+Stripe)", "OVH / web", "E-commerce, paiement"],
        ["Jira Service Mgmt", "Atlassian / cloud", "Support, tickets, reclamations"],
        ["Power BI / MS Project", "Microsoft", "Reporting, pilotage projets"],
        ["DocuFlow / ArchiCAD-Revit", "Torpier / editeurs", "GED, conception CAO/BIM"],
    ], Inches(0.7), Inches(1.5), Inches(12), col_w=[3, 3, 4], fsize=11.5)
    foot(s)

    # 6 Infra
    s = _blank(P); header(s, "1. Modeliser", "Strate technique — composants ↔ applis supportees", "2 pts")
    table(s, [
        ["Composant", "Techno", "Applications / role supportes"],
        ["Serveurs applicatifs", "Win Server 2019/22, Hyper-V", "Applis metiers on-premise"],
        ["Cluster BDD (2 noeuds)", "SQL Server AlwaysOn", "Donnees Dynamics, GesProd, Qualeval"],
        ["Controleurs de domaine", "AD DS (siege + replica Arras)", "Identites & acces, imprimantes"],
        ["Serveur de fichiers / sauvegarde", "DFS / Veeam", "GED, partages, PRA"],
        ["Serveur RDS / Proxy-securite", "Win Server 2019", "Acces distants, filtrage"],
        ["Reseau & interconnexion", "Cisco Catalyst/ISR, VLAN 10-60", "Segmentation par service"],
        ["Interconnexion sites", "VPN IPSec Fortinet, FortiGate 200E", "Nanterre-Arras-Finlande"],
    ], Inches(0.7), Inches(1.5), Inches(12), col_w=[3, 3.2, 4], fsize=11.5)
    foot(s)

    # 7 axes
    s = _blank(P); header(s, "1. Modeliser", "5 axes d'amelioration du SI", "1 pt")
    table(s, [
        ["#", "Axe", "Contenu", "Indicateur"],
        ["1", "Applicatif / urbanisation", "Casser les silos CSV (API/ESB), urbaniser", "Interfaces manuelles supprimees"],
        ["2", "Reseau / interconnexion & resilience", "Fiabiliser/securiser le multi-sites, PRA", "Taux de dispo inter-sites"],
        ["3", "Stockage / donnees & BI", "Qualite donnees, sauvegarde/PRA, Power BI", "RTO / RPO"],
        ["4", "Hebergement / cloud hybride", "Exchange->M365, scalabilite e-commerce", "Cout/perf, elasticite"],
        ["5", "Acces utilisateurs / securite", "IAM, RGPD, acces distants (RDS, MFA)", "Incidents de securite"],
    ], Inches(0.55), Inches(1.55), Inches(12.3), col_w=[0.4, 2.6, 4.2, 2.6], fsize=12)
    _, tf = textbox(s, Inches(0.55), Inches(6.05), Inches(12.3), Inches(0.7))
    set_run(tf.paragraphs[0].add_run(),
            "Liste UNIQUE reprise a l'identique dans la soutenance individuelle (coherence collectif <-> individuel).",
            12.5, TEAL, italic=True, bold=True)
    foot(s)

    # 8 referentiels
    s = _blank(P); header(s, "1. Modeliser", "Comparaison de 3 referentiels de gouvernance", "2 pts")
    table(s, [
        ["Referentiel", "Role", "Interet pour Torpier", "Limite"],
        ["ITIL (retenu)", "Gestion des services / run", "Fiabilise support & exploitation multi-sites", "Peu de gouvernance/projet"],
        ["COBIT", "Gouvernance SI", "Alignement strategie, KPI, RGPD", "Lourd pour une PME"],
        ["Scrum / Agile", "Build / projets dev", "Dev internes (PrestaShop, GesProd), PMO", "Ne couvre pas le run"],
    ], Inches(0.55), Inches(1.6), Inches(12.3), col_w=[1.7, 2.2, 3.6, 2.2], fsize=12.5)
    box = band(s, LIGHT, Inches(0.55), Inches(4.5), Inches(12.3), Inches(1.7))
    _, tf = textbox(s, Inches(0.8), Inches(4.65), Inches(11.8), Inches(1.5))
    set_run(tf.paragraphs[0].add_run(), "Choix : ITIL pour le RUN + Scrum/Agile pour le BUILD, COBIT en chapeau gouvernance.",
            15, NAVY, bold=True)
    p = tf.add_paragraph(); p.space_before = Pt(6)
    set_run(p.add_run(), "→ ITIL est retenu comme fil conducteur ; ce choix est repris tel quel dans la soutenance individuelle.",
            13, GREY, italic=True)
    foot(s)

    # 9 organisation equipes
    s = _blank(P); header(s, "2. Encadrer les equipes internes et/ou externes",
                          "Organisation des equipes SI", "2 pts")
    bullets(s, [
        (0, "Organisation hierarchique + transversale pilotee par le DSI (L. Besnier)", True),
        (1, "Pole RUN / exploitation (ITIL) : A. Duchamp (infra & securite), admins systeme/reseau"),
        (1, "Pole BUILD / applicatif : L. Vigne (applicatif), dev full-stack, alternant, freelances"),
        (1, "Cellule PMO transversale : T. Valentia (PMP/Scrum Master) — uniformise les pratiques"),
        (0, "Style de management adapte (cours CESI : Posture du Manager)", True),
        (1, "Participatif / collaboratif : equipe experte et autonome"),
        (1, "Touches directives sur securite & conformite (coherent avec ITIL)"),
        (1, "Organisation type 'orange->participative' (Mintzberg / Laloux)"),
    ], y=Inches(1.5))
    foot(s)

    # 10 RACI + outils
    s = _blank(P); header(s, "2. Encadrer", "Matrice RACI & outils de pilotage", "RACI 2 + outils 2")
    table(s, [
        ["Activite / projet", "DSI", "Run (infra)", "Build (appli)", "PMO", "Direction"],
        ["Urbanisation / API", "A", "C", "R", "C", "I"],
        ["Interconnexion & PRA", "A", "R", "C", "C", "I"],
        ["Mise en place ITIL", "A", "R", "C", "C", "I"],
        ["Conduite du changement", "A", "C", "C", "R", "C"],
    ], Inches(0.55), Inches(1.5), Inches(12.3),
       col_w=[3.4, 1, 1.4, 1.4, 1, 1.4], fsize=12)
    _, tf = textbox(s, Inches(0.55), Inches(4.4), Inches(12.3), Inches(0.4))
    set_run(tf.paragraphs[0].add_run(), "R=Realise  A=Approuve/responsable  C=Consulte  I=Informe", 11, GREY, italic=True)
    bullets(s, [
        (0, "Outils de SUIVI des projets de transformation : Jira (tickets/Kanban), MS Project, Power BI (KPI)", True),
        (0, "Outils de PILOTAGE du changement : Teams (communication), enquetes d'adhesion, indicateurs ITIL", True),
    ], y=Inches(5.05), h=Inches(1.6), size=15)
    foot(s)

    # 11 pitch changement
    s = _blank(P); header(s, "3. Conduire le changement induit par les projets de la DSI",
                          "Pitch : adopter une conduite du changement", "2 pts")
    bullets(s, [
        (0, "Sans conduite du changement, la transformation echoue : l'adhesion se construit", True),
        (1, "Methode structurante : ITIL + ADKAR (ou 8 etapes de Kotter)"),
        (1, "Leviers : sens (Mission-Vision-Valeurs), implication, reconnaissance, exemplarite du manager"),
        (1, "Communication continue, formation, quick wins visibles (ex. fin des exports CSV manuels)"),
        (0, "Argument direction : moins d'incidents, plus d'agilite, ROI mesurable", True),
        (1, "KPI : disponibilite inter-sites, delai de resolution, satisfaction utilisateurs"),
    ], y=Inches(1.5))
    foot(s)

    # 12 bilan groupe
    s = _blank(P); header(s, "Conclusion", "Bilan du groupe", "1 pt")
    bullets(s, [
        (0, "Points forts du groupe", True),
        (1, "Complementarite des roles, production iterative du support, ancrage dans le contexte Torpier"),
        (0, "Axes d'amelioration du groupe", True),
        (1, "Approfondir le chiffrage des indicateurs, mieux repartir le temps de parole en soutenance"),
        (0, "Prochaine etape : approfondir l'axe 'Optimisation du SI' (soutenance individuelle)", True),
    ], y=Inches(1.6))
    foot(s)

    # 13 merci
    s = title_slide(P, "MERCI", "Questions & echanges",
                    ["Place a la discussion (15 min)"],
                    "Projet collaboratif — cas Torpier")
    foot(s)

    P.save(path)
    print("OK collectif :", path, "-", len(P.slides.__iter__.__self__._sldIdLst), "slides")


# =========================================================================
#  PRESENTATION INDIVIDUELLE
# =========================================================================

def build_individual(path):
    P = new_prez()
    FOOT = "BLOC4 · Manager les equipes et la transformation du SI · Soutenance individuelle (cas Torpier)"
    total = 13
    i = [0]

    def foot(s):
        i[0] += 1
        footer(s, i[0], total, FOOT)

    s = title_slide(P, "SOUTENANCE INDIVIDUELLE — RESPONSABLE SI",
                    "Trajectoire du SI Torpier",
                    ["Aligner le SI sur la strategie & creer de la valeur",
                     "Focus : Optimisation du SI — referentiel ITIL · 20 min + 15 min d'echange"],
                    "Cas entreprise TORPIER — CESI")
    foot(s)

    # accroche
    s = _blank(P); header(s, "Introduction", "Ma mission : aligner le SI sur la strategie")
    bullets(s, [
        (0, "Posture : nouveau responsable SI de Torpier", True),
        (0, "Promesse a la direction : accroitre la valeur ajoutee du SI", True),
        (1, "Garantir que la trajectoire du SI suive celle de l'entreprise"),
        (1, "Fiabiliser, urbaniser et securiser un SI devenu multi-sites"),
        (0, "Fil conducteur : le POURQUOI avant le COMMENT", True),
        (1, "Strategie 2026-2029 : diversification, developpement durable, industrie 4.0, optimisation du SI"),
    ], y=Inches(1.6))
    foot(s)

    # entreprise + processus
    s = _blank(P); header(s, "1. Modeliser les flux metiers et ressources techniques",
                          "Entreprise, activite & processus", "3 pts")
    bullets(s, [
        (0, "Torpier : conception, fabrication et vente de structures/mobiliers bois (espaces publics/prives)", True),
        (1, "Sites : Direction Nanterre · Usine Arras (conception, fabrication) · Usine Finlande"),
        (1, "Transformation : e-commerce, expansion EU, digitalisation, teletravail"),
        (0, "Processus metiers & workflows des flux d'information", True),
        (1, "Commercial/relation client — Production — Financier/administratif — Gouvernance — Supports IT"),
        (1, "Flux cles : commande (Dynamics) -> appro/fabrication (GesProd/Qualeval) -> facturation (SAGE)"),
        (0, "Strategie 2026-2029 (4 axes) : diversification · durable · industrie 4.0 · optimisation du SI", True),
    ], y=Inches(1.45), size=16)
    foot(s)

    # SI graphique
    s = _blank(P); header(s, "1. Modeliser", "SI existant — les 4 strates (cartographie)", "4 pts")
    strates = [
        ("PROCESSUS METIERS", TEAL, "Commercial · Production · Finance/RH · Gouvernance · Supports IT"),
        ("FONCTIONS / APPLICATIONS", RGBColor(0x2E,0x7D,0x46),
         "Dynamics · GesProd · SAGE · Qualeval · PrestaShop · Jira · Power BI · DocuFlow"),
        ("DONNEES & INTEGRATION", ORANGE,
         "SQL Server AlwaysOn · echanges CSV (a remplacer) · GED · web services"),
        ("INFRASTRUCTURE & RESEAU", NAVY,
         "Hyper-V · AD DS · RDS · VLAN 10-60 · VPN Fortinet (Nanterre-Arras-Finlande)"),
    ]
    y = Inches(1.45)
    for lbl, col, txt in strates:
        bar = band(s, col, Inches(0.7), y, Inches(3.6), Inches(1.15))
        bar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        bar.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        set_run(bar.text_frame.paragraphs[0].add_run(), lbl, 13, WHITE, bold=True)
        box = band(s, LIGHT, Inches(4.45), y, Inches(8.15), Inches(1.15))
        box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        _, tf = textbox(s, Inches(4.65), y, Inches(7.8), Inches(1.15), anchor=MSO_ANCHOR.MIDDLE)
        set_run(tf.paragraphs[0].add_run(), txt, 13.5, DARK)
        y += Inches(1.32)
    _, tf = textbox(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.4))
    set_run(tf.paragraphs[0].add_run(),
            "A presenter sous forme de schema d'architecture + topologie reseau commentes.", 11.5, GREY, italic=True)
    foot(s)

    # 5 axes
    s = _blank(P); header(s, "1. Modeliser", "5 axes d'amelioration des flux du SI", "2 pts")
    table(s, [
        ["#", "Axe (flux)", "Amelioration", "Indicateur de valeur"],
        ["1", "Applicatif", "Casser les silos CSV via API/ESB, urbaniser", "Interfaces manuelles supprimees"],
        ["2", "Reseau", "Fiabiliser/securiser l'interconnexion, PRA", "Taux de dispo inter-sites"],
        ["3", "Stockage / donnees", "Qualite donnees, sauvegarde/PRA, BI", "RTO / RPO"],
        ["4", "Hebergement", "Cloud hybride (M365), scalabilite e-commerce", "Cout/perf, elasticite"],
        ["5", "Acces utilisateurs", "IAM, RGPD, acces distants (RDS, MFA)", "Incidents de securite"],
    ], Inches(0.55), Inches(1.55), Inches(12.3), col_w=[0.4, 2.2, 4.4, 2.8], fsize=12)
    _, tf = textbox(s, Inches(0.55), Inches(6.05), Inches(12.3), Inches(0.6))
    set_run(tf.paragraphs[0].add_run(), "Meme liste que la soutenance collective — focus sur les axes 1 & 2.",
            12.5, TEAL, italic=True, bold=True)
    foot(s)

    # focus
    s = _blank(P); header(s, "Focus", "Axe retenu : Optimisation du SI")
    bullets(s, [
        (0, "Pourquoi cet axe ? Il conditionne tous les autres (resilience, agilite, interconnexion)", True),
        (1, "Incarne par les axes 1 & 2 : urbanisation applicative + interconnexion/resilience"),
        (0, "Referentiel : ITIL (choix issu de la comparaison faite en collectif)", True),
        (1, "Gestion des changements, des incidents, des problemes, des niveaux de service"),
        (1, "Fiabilise le RUN multi-sites aujourd'hui informel (Jira sans process)"),
        (0, "Benefices direction : moins d'incidents, SI agile, valeur mesurable", True),
    ], y=Inches(1.6))
    foot(s)

    # 5 projets
    s = _blank(P); header(s, "1. Modeliser", "5 projets operationnels (axe Optimisation du SI)")
    table(s, [
        ["Projet", "Objectif", "Referentiel / process"],
        ["Urbanisation & API GesProd-Dynamics", "Supprimer les echanges CSV", "ITIL Change + Scrum (dev)"],
        ["Fiabilisation interconnexion & PRA", "Resilience Nanterre-Arras-Finlande", "ITIL Service Level / Continuity"],
        ["Mise en place du centre de services ITIL", "Cadrer support & incidents (Jira)", "ITIL Incident/Problem"],
        ["Migration Exchange -> M365", "Messagerie unifiee, cloud hybride", "ITIL Change + projet"],
        ["IAM & securisation des acces", "RGPD, MFA, acces distants", "ITIL Access Mgmt"],
    ], Inches(0.55), Inches(1.6), Inches(12.3), col_w=[3.4, 3.4, 2.6], fsize=12)
    foot(s)

    # equipe + organigramme + management
    s = _blank(P); header(s, "2. Encadrer les equipes internes et/ou externes",
                          "Equipe projet, organigramme & management", "1 + 3 pts")
    bullets(s, [
        (0, "Responsabilites alignees aux competences & aspirations des personas Torpier", True),
        (1, "A. Duchamp (infra/securite) -> interconnexion, PRA, IAM"),
        (1, "L. Vigne (applicatif) -> urbanisation/API, migration M365"),
        (1, "T. Valentia (PMO/Scrum) -> pilotage projets, coordination agile"),
        (1, "Admins systeme/reseau + dev full-stack + alternant ; freelances en renfort"),
        (0, "Type de management : participatif/collaboratif, directif sur securite (coherent ITIL)", True),
        (0, "Outils de suivi : Jira (Kanban/tickets), MS Project (planning), Power BI (KPI)", True),
    ], y=Inches(1.5), size=15)
    foot(s)

    # externes + plan com
    s = _blank(P); header(s, "2. Encadrer", "Fournisseurs/sous-traitants & plan de communication", "2 pts")
    bullets(s, [
        (0, "Besoins externes justifies", True),
        (1, "Freelances dev (absorber la charge), integrateurs M365/ITSM, expertise securite"),
        (0, "Plan de communication prestataires", True),
        (1, "Actions : kick-off, comites de suivi, recette ; Ressources : referent DSI + PMO"),
        (1, "Canaux : Teams (echanges), Jira (tickets partages), e-mail (officiel)"),
        (1, "Indicateurs : respect des delais/SLA, qualite des livrables, satisfaction"),
    ], y=Inches(1.5))
    foot(s)

    # outils collaboratifs
    s = _blank(P); header(s, "2. Encadrer", "Outils collaboratifs retenus")
    table(s, [
        ["Outil", "Usage relation externes", "Usage pilotage interne (KPI)"],
        ["Microsoft Teams", "Echanges fluides avec prestataires", "Rituels d'equipe, communication"],
        ["Jira Service Mgmt", "Tickets partages, process ITIL", "Suivi incidents/changements"],
        ["Power BI", "Reporting partage", "Tableaux de bord de performance"],
    ], Inches(0.7), Inches(1.7), Inches(12), col_w=[2.5, 4, 4], fsize=13)
    foot(s)

    # conduite du changement
    s = _blank(P); header(s, "3. Conduire le changement induit par les projets de la DSI",
                          "Conduite du changement", "2 + 1 pts")
    bullets(s, [
        (0, "Leviers d'engagement : sens (Mission-Vision-Valeurs), implication, reconnaissance, quick wins", True),
        (0, "Methode : ADKAR (ou Kotter) couplee a ITIL", True),
        (1, "Awareness -> Desire -> Knowledge -> Ability -> Reinforcement"),
        (0, "Arguments d'adhesion : moins d'incidents, agilite accrue, ROI mesurable", True),
        (0, "Acteurs cles impactes & communication adaptee", True),
        (1, "Equipe DSI, utilisateurs metiers (commerciaux, production, Finlande), direction, prestataires"),
        (1, "Communication ciblee par partie prenante (ateliers, formations, comites, newsletters)"),
    ], y=Inches(1.5), size=15)
    foot(s)

    s = title_slide(P, "MERCI", "Questions & echanges",
                    ["Place a la discussion (15 min)"],
                    "Soutenance individuelle — cas Torpier")
    foot(s)

    P.save(path)
    print("OK individuel :", path)


if __name__ == "__main__":
    build_collective("Soutenance-COLLECTIVE-Torpier.pptx")
    build_individual("Soutenance-INDIVIDUELLE-Torpier.pptx")
    print("Termine.")
