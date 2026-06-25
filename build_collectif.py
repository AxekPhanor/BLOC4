"""Génère 'BLOC4 collectif.pptx' depuis Template.pptx — idempotent (regénère tout).

Principe : on REUTILISE les mises en page déjà dessinées du template (SWOT, matrices,
timelines, tableaux...) en les personnalisant en place, puis on réordonne les slides
dans l'ordre de la présentation et on supprime les slides non utilisées.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import _design as d

SRC = "Template.pptx"
OUT = "BLOC4 collectif.pptx"

# Catalogue des slides réutilisables du template (index 0-based)
TPL = {
    "title": 0,
    "phases_timeline": 3,    # 4 phases en zigzag
    "matrix_3x3": 4,         # matrice 3x3 (axes High/Medium/Low)
    "decision_tree": 5,      # arbre Question/Yes/No -> résultats
    "donut": 6,
    "worldmap_table": 7,
    "swot": 8,               # SWOT S/W/O/T
    "venn": 9,               # 3 cercles
    "timeline_steps": 14,    # timeline verticale 4 étapes
    "table_3col": 15,        # tableau 3 colonnes x 3 lignes
}


# ----------------------------------------------------------------------------- utils
def reorder_and_prune(prs, ordered):
    """Réordonne les slides selon `ordered` (liste d'objets slide) et supprime le reste."""
    sldIdLst = prs.slides._sldIdLst
    els = list(sldIdLst)
    by_id = {sl.slide_id: el for el, sl in zip(els, list(prs.slides))}
    ordered_ids = [sl.slide_id for sl in ordered]  # capturer avant détachement
    for el in els:
        sldIdLst.remove(el)
    keep = set()
    for sid in ordered_ids:
        el = by_id[sid]
        sldIdLst.append(el)
        keep.add(el.get(qn("r:id")))
    for el in els:
        rid = el.get(qn("r:id"))
        if rid not in keep:
            prs.part.drop_rel(rid)


def _set_run_text(tb, text, size=None):
    """Remplace le texte en gardant la mise en forme du 1er run (taille optionnelle)."""
    p = tb.text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for r in p.runs[1:]:
            r._r.getparent().remove(r._r)
        if size is not None:
            p.runs[0].font.size = Pt(size)
    else:
        p.text = text


def reuse_clean(slide, title):
    """Nettoie une slide du template (titre générique + crédit Slidesgo) et pose le bandeau orange."""
    for sh in list(slide.shapes):
        if sh.has_text_frame:
            t = sh.text_frame.text.lower()
            if "slidesgo" in t or "business infographics" in t:
                sh._element.getparent().remove(sh._element)
    d.header(slide, title)


def fill_swot(slide, forces, faiblesses, opportunites, menaces):
    """Personnalise le SWOT du template (titres FR + descriptions)."""
    data = {
        "strengths": ("Forces", forces),
        "weaknesses": ("Faiblesses", faiblesses),
        "opportunities": ("Opportunités", opportunites),
        "threats": ("Menaces", menaces),
    }
    for grp in slide.shapes:
        if grp.shape_type != 6:  # GROUP
            continue
        inner_groups = [c for c in grp.shapes if c.shape_type == 6]
        if not inner_groups:
            continue
        tboxes = [sub for ig in inner_groups for sub in ig.shapes
                  if sub.has_text_frame and sub.text_frame.text.strip()]
        head_tb = next((tb for tb in tboxes if tb.text_frame.text.strip().lower() in data), None)
        if head_tb is None:
            continue
        desc_tb = next((tb for tb in tboxes if tb is not head_tb), None)
        newhead, newdesc = data[head_tb.text_frame.text.strip().lower()]
        _set_run_text(head_tb, newhead)
        if desc_tb is not None:
            _set_run_text(desc_tb, newdesc, size=10)


# ----------------------------------------------------------------------------- slides
def build_d1(s):
    """D1 — Titre (slide TITLE du template)."""
    s.placeholders[0].text = "Manager les équipes et la transformation du SI"
    sub = s.placeholders[1]
    sub.text = "Groupe Torpier — Projet collaboratif DSI"
    sub.text_frame.add_paragraph().text = "[Membres du groupe]"
    for sh in list(s.shapes):
        if sh.has_text_frame and "slidesgo" in sh.text_frame.text.lower():
            sh._element.getparent().remove(sh._element)
    d.notes(s, """
Slide d'ouverture — présentation de l'équipe et du cadre.
• Se présenter : « Nous sommes l'équipe projet de la DSI du Groupe Torpier, constituée en groupe de travail à la demande du Directeur SI. »
• Rappeler la mission : réaliser une étude préliminaire d'évolution du SI (ébauche de schéma directeur) pour le garder aligné avec la stratégie de l'entreprise.
• Annoncer le plan : 1) analyse de l'existant, 2) gouvernance & organisation des équipes, 3) outils & conduite du changement, 4) organisation du groupe & bilan.
• Durée : 20 min de présentation puis 15 min d'échange. Objectif : convaincre la direction par l'argumentation.
""")


def build_d2(prs):
    """D2 — Activités & stratégie (slide construite, 2 colonnes)."""
    s = d.blank_slide(prs)
    d.header(s, "Le Groupe Torpier — activités & stratégie")
    LX, RX, W = Inches(0.5), Inches(5.15), Inches(4.35)
    CY, CH = Inches(1.45), Inches(3.5)
    d.card(s, LX, CY, W, CH, fill=d.CREME)
    d.card(s, RX, CY, W, CH, fill=d.CREME)
    d.chip(s, LX + Inches(0.18), CY + Inches(0.18), Inches(3.1), "Le Groupe & ses activités", color=d.ORANGE)
    d.chip(s, RX + Inches(0.18), CY + Inches(0.18), Inches(3.0), "Stratégie 2026-2029", color=d.ANTHRA)

    tb, tf = d.textbox(s, LX + Inches(0.18), CY + Inches(0.62), W - Inches(0.36), CH - Inches(0.78))
    d.set_para(tf.paragraphs[0], "Spécialiste des aménagements d'extérieur, de la conception/R&D à la production et la commercialisation. Deux filiales :", size=11.5, color=d.ANTHRA)
    d.add_para(tf, "TORPIER SAS — structures & mobilier extérieurs en bois (charpente, pergolas, abris-bus)", size=11.5, color=d.ORANGE, bold=True, space_before=7)
    d.add_para(tf, "POTIER SARL — mobilier urbain & luminaire extérieurs en aluminium (leader français)", size=11.5, color=d.ORANGE, bold=True, space_before=4)
    d.add_para(tf, "•  Sites : siège de Nanterre · usine d'Arras · usine en Finlande · bureaux commerciaux en Allemagne, Italie et Espagne", size=11.5, color=d.ANTHRA, space_before=7)
    d.add_para(tf, "•  +30 % de CA (2023-2025), virage e-commerce, digitalisation des processus & télétravail", size=11.5, color=d.ANTHRA, space_before=5)

    tb2, tf2 = d.textbox(s, RX + Inches(0.18), CY + Inches(0.62), W - Inches(0.36), CH - Inches(0.78))
    d.set_para(tf2.paragraphs[0], "4 axes définis par la Direction (S. Crélin) :", size=11.5, color=d.ANTHRA, bold=True)
    d.add_para(tf2, "1.  Diversification des cibles clients (privé, promoteurs immobiliers)", size=12, space_before=9)
    d.add_para(tf2, "2.  Développement durable (traçabilité bois, certification fournisseurs)", size=12, space_before=7)
    d.add_para(tf2, "3.  Industrie 4.0 (IoT, maintenance prédictive, automatisation)", size=12, space_before=7)
    d.add_para(tf2, "4.  Optimisation du SI — interconnexion des sites, résilience & agilité", size=12, color=d.ORANGE, bold=True, space_before=7)
    d.add_para(tf2, "→ périmètre de la DSI : notre point de départ", size=11, color=d.ORANGE, bold=True, space_before=2)
    d.footer(s, "Activités & stratégie")
    d.notes(s, """
Poser le décor — qui est Torpier et où veut aller la direction.
• Le Groupe = 2 filiales complémentaires sur l'aménagement d'extérieur : TORPIER SAS (bois) et POTIER SARL (mobilier urbain & luminaire en aluminium, leader français). Conception/R&D, production et commercialisation intégrées.
• Implantations : siège à Nanterre, usine d'Arras, usine en Finlande (rachat récent), bureaux commerciaux en Allemagne, Italie et Espagne. → un SI multi-sites et multi-filiales.
• Dynamique : +30 % de CA sur 2023-2025, virage e-commerce, digitalisation et télétravail qui sollicitent fortement le SI.
• Stratégie 2026-2029 (Sophie Crélin, DG) : 4 axes. Insister sur le « pourquoi » : les 3 premiers axes (clients, durable, industrie 4.0) reposent tous sur un SI solide.
• C'est l'axe 4 « Optimisation du SI » qui est notre périmètre : interconnexion des sites, résilience et agilité. → fil rouge de toute la présentation.
• Anticiper Q/R : la croissance par rachats successifs a créé une hétérogénéité applicative entre filiales (silos) — on y revient dans le SWOT et l'analyse du SI.
• Note interne : Pologne (plan SDSI 2019) = même opération d'expansion que la Finlande (réalisée, contexte 2025) ; on retient la Finlande.
""")
    return s


def build_d3(s):
    """D3 — SWOT de l'entreprise (slide SWOT du template personnalisée)."""
    reuse_clean(s, "SWOT de l'entreprise")
    fill_swot(
        s,
        "2 filiales complémentaires (bois + aluminium/luminaire), chaîne intégrée, forte croissance.",
        "Filiales en silos, SI hétérogène, processus manuels (Excel/mail), support sans indicateurs.",
        "Marché privé & promoteurs, demande durable, industrie 4.0, expansion européenne.",
        "Approvisionnement tendu, pression réglementaire (RGPD), risques cyber & continuité.",
    )
    d.notes(s, """
Diagnostic stratégique de l'entreprise (pas encore du SI) — il cadre nos propositions.
• FORCES : complémentarité des 2 filiales (bois + aluminium/luminaire) → offre large ; chaîne de valeur intégrée ; croissance avec e-commerce en place.
• FAIBLESSES : la croissance par rachats a laissé des filiales en silos, un SI hétérogène, trop de processus manuels (Excel/mail, cf. témoignages) et un support IT sans indicateurs.
• OPPORTUNITÉS : nouveaux marchés (privé, promoteurs), durable et personnalisation, industrie 4.0 et expansion européenne.
• MENACES : approvisionnement tendu (vécu 2020-2022), pression réglementaire, risques cyber/continuité non couverts (pas de PRA).
• Transition : plusieurs faiblesses/menaces sont adressables par le SI → on passe à la cartographie de l'existant.
""")
    return s


GREY = RGBColor(0x6B, 0x6B, 0x6B)
TINT = RGBColor(0xFC, 0xE9, 0xDF)  # orange très clair (surlignage divergence)


def _cell(cell, text, size=9.5, color=d.ANTHRA, bold=False, fill=d.BLANC, align=PP_ALIGN.LEFT):
    """Style une cellule de tableau (fond, marges, police)."""
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.05)
    cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
    cell.text_frame.word_wrap = True
    d.set_para(cell.text_frame.paragraphs[0], text, size=size, color=color, bold=bold, align=align)


def build_d5(prs):
    """D5 — Tableau Applications du SI ↔ fonctions métiers (par filiale)."""
    s = d.blank_slide(prs)
    d.header(s, "Applications du SI ↔ fonctions métiers",
             "Deux filiales, des applications hétérogènes pour les mêmes fonctions")

    data = [
        ("Relation client (CRM)",        "Microsoft Dynamics", "Zoho"),
        ("Production",                   "GesProd",            "GesProd"),
        ("Gestion économique & financière", "SAGE Compta",     "ISA Compta"),
        ("Paie & ressources humaines",   "SAGE Paie",          "ISA Compta"),
        ("Gestion des temps",            "OCTIME",             "— (Excel)"),
        ("Qualité",                      "Qualéval",           "—"),
        ("Portail décisionnel (BI)",     "SAP BusinessObjects","—"),
        ("Reporting d'activité",         "I-MAGE",             "I-MAGE"),
        ("Gestion de projets",           "SuiteProG",          "Excel"),
    ]

    X, Y, W = Inches(0.5), Inches(1.4), Inches(9.0)
    RH = Inches(0.32)
    n = len(data) + 1
    gfx = s.shapes.add_table(n, 3, X, Y, W, RH * n)
    table = gfx.table
    table.first_row = False
    table.horz_banding = False
    table.columns[0].width = Inches(3.4)
    table.columns[1].width = Inches(2.8)
    table.columns[2].width = Inches(2.8)
    for r in table.rows:
        r.height = RH

    for j, h in enumerate(("Fonction métier", "TORPIER SAS", "POTIER SARL")):
        _cell(table.cell(0, j), h, size=11, color=d.BLANC, bold=True,
              fill=d.ANTHRA if j == 0 else d.ORANGE,
              align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

    for i, (fn, a, b) in enumerate(data, start=1):
        app_fill = TINT if a != b else d.BLANC  # divergence (ou manque) = silo
        _cell(table.cell(i, 0), fn, size=9.5, color=d.ANTHRA, bold=True, fill=d.CREME)
        _cell(table.cell(i, 1), a, size=9.5, color=d.ANTHRA, fill=app_fill, align=PP_ALIGN.CENTER)
        _cell(table.cell(i, 2), b, size=9.5, color=d.ANTHRA, fill=app_fill, align=PP_ALIGN.CENTER)

    tb, tf = d.textbox(s, X, Y + RH * n + Inches(0.1), W, Inches(0.35))
    d.set_para(tf.paragraphs[0],
               "Cellules orangées = applications divergentes entre filiales (silos). Seuls GesProd et I-MAGE sont communs → cible : ERP unique (projets P7–P8) & urbanisation du SI.",
               size=9, color=d.ORANGE, bold=True)

    d.footer(s, "Applis ↔ fonctions")
    d.notes(s, """
Tableau de correspondance applications ↔ fonctions métiers, vu par filiale — il matérialise l'hétérogénéité du SI.
• Lecture : chaque ligne = une fonction métier ; les colonnes = les applications utilisées par TORPIER SAS et par POTIER SARL.
• Constat majeur : pour une même fonction, les deux filiales utilisent des outils différents (Dynamics vs Zoho pour le CRM, SAGE vs ISA pour la finance/paie), héritage des rachats successifs → SI en silos, doublons de licences, données non consolidées.
• Seules GesProd (production) et I-MAGE (reporting) sont communes aux deux filiales.
• Côté POTIER, plusieurs fonctions ne sont pas outillées (temps, qualité, BI) ou reposent sur Excel → manque de fiabilité et de pilotage.
• Conclusion / transition : cette cartographie applicative justifie directement l'axe « Uniformiser & urbaniser le SI » et le projet d'ERP unique (lots P7-P8 du SDSI). Elle prépare le tableau infrastructure ↔ applications (slide suivante).
""")
    return s


def rail(s, x, y, w, h, text, color, tcolor=d.BLANC):
    """Onglet vertical de famille de processus (libellé à gauche d'une bande)."""
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    tf = r.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.03); tf.margin_right = Inches(0.03)
    d.set_para(tf.paragraphs[0], text, size=8.5, color=tcolor, bold=True, align=PP_ALIGN.CENTER)
    return r


def build_d4(prs):
    """D4 — Cartographie des processus métiers (3 familles, tâches & acteurs)."""
    s = d.blank_slide(prs)
    d.header(s, "Cartographie des processus métiers",
             "Outillés par le SI — tâches principales & acteurs")

    RAILX, RAILW = Inches(0.5), Inches(1.05)
    CX, CW = Inches(1.7), Inches(7.8)

    # --- Bande PILOTAGE ---
    py, ph = Inches(1.35), Inches(0.82)
    rail(s, RAILX, py, RAILW, ph, "PILOTAGE", d.ANTHRA)
    d.card(s, CX, py, CW, ph, fill=d.CREME)
    d.chip(s, CX + Inches(0.12), py + Inches(0.24), Inches(2.4), "Pilotage & gouvernance", color=d.ANTHRA)
    tb, tf = d.textbox(s, CX + Inches(2.65), py + Inches(0.08), CW - Inches(2.8), ph - Inches(0.16),
                       anchor=MSO_ANCHOR.MIDDLE)
    d.set_para(tf.paragraphs[0], "Stratégie & schéma directeur SI · pilotage budgétaire · performance (KPI)",
               size=9.5, color=d.ANTHRA)
    d.add_para(tf, "Acteurs : Direction Générale · DSI · directions de filiales",
               size=9, color=d.ORANGE, bold=True, space_before=2)

    # --- Bande RÉALISATION (chaîne de valeur) ---
    ry, rh = Inches(2.28), Inches(2.05)
    rail(s, RAILX, ry, RAILW, rh, "RÉALISATION", d.ORANGE)
    procs = [
        ("Commercial & relation client",
         ["Prospection & devis", "Commande (ADV)", "Vente en ligne · SAV"],
         "Force de vente · ADV · E-commerce"),
        ("Conception & R&D",
         ["Bureau d'études", "Conception CAO / BIM", "Prototypage · nomenclatures"],
         "Bureau d'études · R&D"),
        ("Production & qualité",
         ["Planification (GPAO)", "Fabrication Arras / Finlande", "Contrôle qualité"],
         "Ateliers · Méthodes · Qualité"),
        ("Logistique & livraison",
         ["Gestion des stocks", "Achats & approvisionnement", "Expédition · livraison"],
         "Supply chain · Magasin · Achats"),
    ]
    pw, gap = Inches(1.82), Inches(0.14)
    for i, (t, tasks, act) in enumerate(procs):
        x = CX + (pw + gap) * i
        d.card(s, x, ry, pw, rh, fill=d.CREME)
        d.chip(s, x + Inches(0.08), ry + Inches(0.1), pw - Inches(0.16), t, color=d.ORANGE, h=Inches(0.5))
        tb, tf = d.textbox(s, x + Inches(0.12), ry + Inches(0.7), pw - Inches(0.24), Inches(0.85))
        d.set_para(tf.paragraphs[0], "•  " + tasks[0], size=9, color=d.ANTHRA)
        for tk in tasks[1:]:
            d.add_para(tf, "•  " + tk, size=9, color=d.ANTHRA, space_before=3)
        tb2, tf2 = d.textbox(s, x + Inches(0.12), ry + rh - Inches(0.46), pw - Inches(0.24), Inches(0.42),
                             anchor=MSO_ANCHOR.BOTTOM)
        d.set_para(tf2.paragraphs[0], act, size=8, color=d.ORANGE, bold=True)
        if i < 3:
            cb, cf = d.textbox(s, x + pw - Inches(0.02), ry + Inches(0.78), gap + Inches(0.06), Inches(0.5),
                               anchor=MSO_ANCHOR.MIDDLE)
            d.set_para(cf.paragraphs[0], "›", size=22, color=d.ORANGE, bold=True, align=PP_ALIGN.CENTER)

    # --- Bande SUPPORT ---
    sy, sbh = Inches(4.44), Inches(0.8)
    rail(s, RAILX, sy, RAILW, sbh, "SUPPORT", GREY)
    sup = [
        ("Finance & comptabilité", "DAF · Comptabilité"),
        ("Ressources humaines & paie", "RH · Paie"),
        ("Systèmes d'information", "DSI · Support IT"),
    ]
    sw, sgap = Inches(2.5), Inches(0.15)
    for i, (t, act) in enumerate(sup):
        x = CX + (sw + sgap) * i
        d.card(s, x, sy, sw, sbh, fill=d.CREME)
        d.chip(s, x + Inches(0.1), sy + Inches(0.1), sw - Inches(0.2), t, color=GREY)
        tb, tf = d.textbox(s, x + Inches(0.12), sy + Inches(0.47), sw - Inches(0.24), Inches(0.3))
        d.set_para(tf.paragraphs[0], act, size=8.5, color=d.ANTHRA)

    d.footer(s, "Processus")
    d.notes(s, """
Cartographie macro des processus métiers de Torpier — lecture en 3 familles, tous outillés par le SI.
• PILOTAGE : la Direction Générale (S. Crélin), la DSI et les directions de filiales fixent la stratégie, le schéma directeur SI, le budget et suivent la performance.
• RÉALISATION (cœur de métier) : la chaîne de valeur de l'aménagement extérieur, du commercial à la livraison.
   1. Commercial & relation client : prospection, devis, prise de commande (ADV), vente en ligne et SAV — force de vente, ADV, e-commerce.
   2. Conception & R&D : bureau d'études, conception CAO/BIM, prototypage et nomenclatures.
   3. Production & qualité : planification (GPAO), fabrication sur les usines d'Arras et de Finlande, contrôle qualité.
   4. Logistique & livraison : stocks, achats/approvisionnement, expédition.
• SUPPORT : finance & comptabilité, RH & paie, et la DSI/support IT qui soutiennent l'ensemble.
• Message clé : chaque processus s'appuie sur des applications du SI — c'est ce que détaille le tableau applications ↔ fonctions métiers (slide suivante). Les processus en silos entre filiales justifient l'axe « Uniformiser & urbaniser le SI ».
""")
    return s


def num_badge(s, x, y, n, size=Inches(0.44), color=d.ORANGE):
    """Pastille ronde numérotée (badge orange, chiffre blanc centré)."""
    b = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    b.fill.solid(); b.fill.fore_color.rgb = color
    b.line.fill.background(); b.shadow.inherit = False
    tf = b.text_frame
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    d.set_para(tf.paragraphs[0], str(n), size=16, color=d.BLANC, bold=True, align=PP_ALIGN.CENTER)
    return b


def build_d7(prs):
    """D7 — 5 axes d'amélioration du SI (+ 1 indicateur de valeur par axe)."""
    s = d.blank_slide(prs)
    d.header(s, "5 axes d'amélioration du SI",
             "Cohérents avec la stratégie 2026-2029 et le diagnostic de l'existant")

    axes = [
        ("Uniformiser et urbaniser le SI",
         "Silos applicatifs entre filiales → cible & cartographie communes",
         "% d'applications cartographiées / urbanisées"),
        ("Move to Cloud",
         "Infrastructure vieillissante → hébergement cloud, élasticité & PRA",
         "Disponibilité (uptime) & part des charges migrées"),
        ("Référentiel unique de gouvernance",
         "Pratiques DSI informelles → cadre commun ITIL (projets & services)",
         "% de processus DSI formalisés et outillés"),
        ("Gestion des services et des incidents",
         "Support IT sans suivi → centre de services ITSM & engagements (SLA)",
         "Délai moyen de résolution (MTTR) & respect des SLA"),
        ("Sécurité et conformité",
         "Exposition cyber & RGPD non maîtrisée → sécuriser et mettre en conformité",
         "Conformité RGPD & vulnérabilités critiques traitées"),
    ]

    X, W = Inches(0.5), Inches(9.0)
    Y0, RH, STEP = Inches(1.4), Inches(0.68), Inches(0.78)
    SEPX = Inches(6.25)  # début du bloc indicateur
    for i, (titre, why, ind) in enumerate(axes):
        y = Y0 + STEP * i
        d.card(s, X, y, W, RH, fill=d.CREME)
        num_badge(s, X + Inches(0.13), y + Inches(0.12), i + 1)
        tb, tf = d.textbox(s, X + Inches(0.72), y + Inches(0.05), Inches(5.45), RH - Inches(0.1),
                           anchor=MSO_ANCHOR.MIDDLE)
        d.set_para(tf.paragraphs[0], titre, size=13.5, color=d.ANTHRA, bold=True)
        d.add_para(tf, why, size=9.5, color=d.ANTHRA, space_before=1)
        d.rect(s, SEPX, y + Inches(0.12), Pt(1.5), RH - Inches(0.24), d.GRIS)
        tb2, tf2 = d.textbox(s, SEPX + Inches(0.18), y + Inches(0.05), Inches(2.55), RH - Inches(0.1),
                             anchor=MSO_ANCHOR.MIDDLE)
        d.set_para(tf2.paragraphs[0], "INDICATEUR", size=8, color=d.ORANGE, bold=True)
        d.add_para(tf2, ind, size=9.5, color=d.ANTHRA, space_before=1)

    d.footer(s, "5 axes")
    d.notes(s, """
Les 5 axes d'amélioration du SI — ils découlent du diagnostic de l'existant et servent la stratégie 2026-2029.
• Insister sur la cohérence : chaque axe répond à une faiblesse constatée (le « → » de chaque ligne) et porte un indicateur de valeur pour mesurer le bénéfice et faciliter l'adhésion de la direction.
1. Uniformiser et urbaniser le SI : sortir des silos applicatifs hérités des rachats, poser une cartographie et une cible d'urbanisation communes aux filiales.
2. Move to Cloud : remplacer une infra vieillissante par un hébergement cloud élastique et résilient (avec PRA), prérequis de la continuité et de l'industrie 4.0.
3. Référentiel unique de gouvernance : adopter un cadre commun (ITIL) pour piloter de façon homogène projets et services de la DSI.
4. Gestion des services et des incidents : structurer un centre de services (ITSM) avec des engagements de niveau (SLA) là où le support était informel et sans suivi.
5. Sécurité et conformité : maîtriser l'exposition cyber et la conformité RGPD, aujourd'hui non couvertes.
• Garde-fou de cohérence : ce sont exactement les 5 axes repris dans les livrables individuels ; chaque membre en approfondira un seul.
""")
    return s


def build_d6(prs):
    """D6 — Strate technique (serveurs + réseau) ↔ applications supportées."""
    s = d.blank_slide(prs)
    d.header(s, "Strate technique ↔ applications supportées",
             "Le socle d'infrastructure et de réseau sous le SI applicatif")

    data = [
        ("Serveurs hôtes Hyper-V", "Virtualisation (VM applicatives)",
         "Dynamics, GesProd, SAGE / ISA, Qualéval"),
        ("Cluster SQL Server", "Bases de données",
         "Dynamics, GesProd, SuiteProG, BO / I-MAGE"),
        ("Active Directory + DNS/DHCP", "Annuaire & authentification",
         "Accès et SSO de toutes les applications"),
        ("Serveur de fichiers (partages)", "Stockage bureautique & GED",
         "Documents, Qualéval, fichiers Excel"),
        ("Serveur de messagerie", "Mail & agenda (à mutualiser, P10)",
         "Messagerie, collaboration"),
        ("Sauvegarde Veeam", "Sauvegarde & reprise (PRA)",
         "Continuité de l'ensemble des applications"),
        ("Pare-feu FortiGate + VPN IPSec", "Sécurité périmétrique & accès distant",
         "E-commerce, télétravail, liaisons inter-sites"),
        ("Switchs Cisco + VLAN 3/20/30/50/66/99", "Segmentation réseau (LAN, WiFi, Admin)",
         "Postes, serveurs, WiFi invité"),
        ("Liaison MPLS inter-sites", "WAN entre sites",
         "Relie Nanterre · Arras · Finlande · bureaux UE"),
    ]

    X, Y, W = Inches(0.5), Inches(1.4), Inches(9.0)
    RH = Inches(0.32)
    n = len(data) + 1
    gfx = s.shapes.add_table(n, 3, X, Y, W, RH * n)
    table = gfx.table
    table.first_row = False
    table.horz_banding = False
    table.columns[0].width = Inches(3.0)
    table.columns[1].width = Inches(3.0)
    table.columns[2].width = Inches(3.0)
    for r in table.rows:
        r.height = RH

    for j, h in enumerate(("Composant technique", "Rôle", "Applications / services supportés")):
        _cell(table.cell(0, j), h, size=10.5, color=d.BLANC, bold=True,
              fill=d.ANTHRA if j == 0 else d.ORANGE)

    for i, (comp, role, apps) in enumerate(data, start=1):
        _cell(table.cell(i, 0), comp, size=9, color=d.ANTHRA, bold=True, fill=d.CREME)
        _cell(table.cell(i, 1), role, size=9, color=d.ANTHRA, fill=d.BLANC)
        _cell(table.cell(i, 2), apps, size=9, color=d.ANTHRA, fill=d.BLANC)

    tb, tf = d.textbox(s, X, Y + RH * n + Inches(0.1), W, Inches(0.35))
    d.set_para(tf.paragraphs[0],
               "Socle vieillissant et peu résilient (PRA partiel, messagerie non mutualisée) → justifie l'axe « Move to Cloud » et la sécurisation du SI.",
               size=9, color=d.ORANGE, bold=True)

    d.footer(s, "Infra ↔ applis")
    d.notes(s, """
Tableau infrastructure ↔ applications — la couche basse qui héberge tout le SI applicatif.
• Lecture : chaque composant technique (serveur ou élément réseau) et les applications/services qu'il supporte.
• Virtualisation Hyper-V et cluster SQL Server hébergent le cœur applicatif (ERP, CRM, production).
• Active Directory gère l'authentification et les accès de toutes les applications ; le serveur de fichiers porte la bureautique et la GED.
• Réseau : switchs Cisco segmentés en VLAN (LAN, WiFi invité, MPLS, Admin), pare-feu FortiGate + VPN pour l'accès distant et l'e-commerce, liaison MPLS reliant les sites (Nanterre, Arras, Finlande, bureaux UE).
• Fragilités : PRA seulement partiel via Veeam, messagerie non mutualisée (projet P10), socle vieillissant.
• Transition : ces fragilités justifient l'axe « Move to Cloud » (élasticité, résilience, PRA) et l'axe sécurité ; elles nourrissent les 5 axes d'amélioration.
""")
    return s


def build_d8(prs):
    """D8 — Comparaison des 3 référentiels (ITIL / COBIT / Scrum-Agile)."""
    s = d.blank_slide(prs)
    d.header(s, "Quel référentiel pour piloter la transformation ?",
             "Comparaison ITIL · COBIT · Scrum-Agile et choix pour Torpier")

    refs = [
        ("ITIL", d.ORANGE, "Gestion des services IT (support, incidents, changements, SLA)",
         "Structure un support aujourd'hui informel, met en place des SLA, colle aux projets P9/P10 et à un SI multi-sites en exploitation",
         "Peu orienté développement/agilité ; lourd si appliqué intégralement",
         "RETENU — socle principal", d.ORANGE, d.BLANC),
        ("COBIT", d.ANTHRA, "Gouvernance & audit du SI (alignement, contrôle, conformité)",
         "Cadre de gouvernance et de conformité (RGPD), pilotage par la valeur et les risques pour la direction",
         "Très normatif et orienté audit ; lourd pour une ETI",
         "En appui — gouvernance & conformité", d.ANTHRA, d.BLANC),
        ("Scrum / Agile", GREY, "Gestion itérative de projets et de produits",
         "Idéal pour les projets de dév (CRM, e-commerce, BI) : time-to-market et forte implication des métiers",
         "Ne couvre ni l'exploitation ni la gouvernance globale",
         "En appui — projets de dév", GREY, d.BLANC),
    ]

    CW, gap = Inches(2.92), Inches(0.12)
    Y, CH = Inches(1.4), Inches(3.45)
    labels = ("À quoi ça sert", "Atouts pour Torpier", "Limites")
    for i, (name, col, role, atout, limite, verdict, vcol, vtc) in enumerate(refs):
        x = Inches(0.5) + (CW + gap) * i
        d.card(s, x, Y, CW, CH, fill=d.CREME)
        d.chip(s, x + Inches(0.12), Y + Inches(0.12), CW - Inches(0.24), name, color=col, h=Inches(0.42))
        tb, tf = d.textbox(s, x + Inches(0.16), Y + Inches(0.66), CW - Inches(0.32), CH - Inches(1.25))
        for lab, txt, first in ((labels[0], role, True), (labels[1], atout, False), (labels[2], limite, False)):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            d.set_para(p, lab, size=8.5, color=col, bold=True, space_before=0 if first else 7)
            d.add_para(tf, txt, size=9, color=d.ANTHRA, space_before=1)
        vb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.12), Y + CH - Inches(0.46),
                                CW - Inches(0.24), Inches(0.36))
        vb.fill.solid(); vb.fill.fore_color.rgb = vcol
        vb.line.fill.background(); vb.shadow.inherit = False
        vb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        d.set_para(vb.text_frame.paragraphs[0], verdict, size=9.5, color=vtc, bold=True, align=PP_ALIGN.CENTER)

    tb, tf = d.textbox(s, Inches(0.5), Y + CH + Inches(0.08), Inches(9.0), Inches(0.4))
    d.set_para(tf.paragraphs[0],
               "Choix Torpier : ITIL comme socle (services & support), complété par COBIT (gouvernance / conformité RGPD) et Scrum-Agile (projets de dév).",
               size=10, color=d.ORANGE, bold=True, align=PP_ALIGN.CENTER)

    d.footer(s, "Référentiels")
    d.notes(s, """
Comparaison argumentée des 3 référentiels — critère le plus pondéré de la grille.
• ITIL : référentiel de gestion des services IT. Pour Torpier, il structure un support aujourd'hui informel (cf. témoignages : incidents par mail/téléphone, priorisation « au plus fort »), instaure des SLA et une CMDB. Limite : peu orienté dév/agilité. → RETENU comme socle.
• COBIT : gouvernance et audit du SI, pilotage par la valeur et les risques, conformité. Très utile pour la direction et le RGPD, mais normatif et lourd. → en appui sur la gouvernance.
• Scrum/Agile : gestion itérative de projets/produits. Parfait pour les projets de dév (CRM, e-commerce, BI) et l'implication des métiers, mais ne couvre pas l'exploitation. → en appui sur les projets.
• Message : on ne choisit pas « un contre les autres » — on combine, avec ITIL comme colonne vertébrale. Cohérent avec le focus individuel ITIL.
""")
    return s


def orgbox(s, x, y, w, h, title, sub=None, fill=d.CREME, titlecolor=d.ANTHRA, subcolor=d.ANTHRA):
    """Boîte d'organigramme (titre + sous-titre centrés)."""
    d.card(s, x, y, w, h, fill=fill)
    tb, tf = d.textbox(s, x + Inches(0.05), y + Inches(0.03), w - Inches(0.1), h - Inches(0.06),
                       anchor=MSO_ANCHOR.MIDDLE)
    d.set_para(tf.paragraphs[0], title, size=10, color=titlecolor, bold=True, align=PP_ALIGN.CENTER)
    if sub:
        d.add_para(tf, sub, size=8, color=subcolor, align=PP_ALIGN.CENTER, space_before=2)


def build_d9(prs):
    """D9 — Organisation des équipes SI (hiérarchique + transversale)."""
    s = d.blank_slide(prs)
    d.header(s, "Organisation des équipes SI",
             "Une ligne hiérarchique (DSI) activée en mode projet transversal")

    # Boîte sommet : DSI / RSSI
    orgbox(s, Inches(3.9), Inches(1.4), Inches(2.2), Inches(0.62),
           "DSI / RSSI", "Ludovic Besnier", fill=d.ANTHRA, titlecolor=d.BLANC, subcolor=d.GRIS)
    # Bus + connecteurs
    d.rect(s, Inches(5.0) - Pt(0.75), Inches(2.02), Pt(1.5), Inches(0.3), d.GRIS)  # vertical
    d.rect(s, Inches(1.5), Inches(2.32), Inches(7.0), Pt(1.5), d.GRIS)             # horizontal bus

    branches = [
        ("Chefs de projet SI", "F. Cerpiat · G. Salas"),
        ("Développement", "W. Roux · M. Tamiari"),
        ("Exploitation", "C. Legrand · A. Ripaud"),
        ("Support utilisateur", "A. Pilon"),
    ]
    bw = Inches(1.95)
    bx0 = Inches(0.62)
    step = Inches(2.12)
    for i, (t, sub) in enumerate(branches):
        x = bx0 + step * i
        cx = x + bw / 2
        d.rect(s, cx - Pt(0.75), Inches(2.32), Pt(1.5), Inches(0.28), d.GRIS)  # descente
        orgbox(s, x, Inches(2.6), bw, Inches(0.92), t, sub, fill=d.CREME)

    # Bande transversale (mode projet)
    by = Inches(3.95)
    band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), by, Inches(9.0), Inches(1.0))
    band.fill.solid(); band.fill.fore_color.rgb = d.CREME
    band.line.color.rgb = d.ORANGE; band.line.width = Pt(1.25); band.shadow.inherit = False
    d.chip(s, Inches(0.7), by + Inches(0.16), Inches(3.0), "Gouvernance transversale — mode projet", color=d.ORANGE)
    tb, tf = d.textbox(s, Inches(3.9), by + Inches(0.1), Inches(5.4), Inches(0.8), anchor=MSO_ANCHOR.MIDDLE)
    d.set_para(tf.paragraphs[0],
               "Comités de pilotage projet · référents métiers des filiales (TORPIER & POTIER) · RSSI en sécurité transverse · prestataires externes",
               size=9.5, color=d.ANTHRA)

    tb2, tf2 = d.textbox(s, Inches(0.5), Inches(5.08), Inches(9.0), Inches(0.32))
    d.set_para(tf2.paragraphs[0],
               "Vertical = lien hiérarchique (autorité) · Horizontal = coopération transversale par projet (cartographie des compétences).",
               size=9, color=d.ORANGE, bold=True, align=PP_ALIGN.CENTER)

    d.footer(s, "Équipes SI")
    d.notes(s, """
Organisation des équipes SI — articuler la hiérarchie et le mode projet.
• Ligne hiérarchique : le DSI (Ludovic Besnier), également RSSI, encadre 4 pôles — chefs de projet (Cerpiat, Salas), développement (Roux, Tamiari), exploitation (Legrand, Ripaud) et support utilisateur (Pilon).
• Mode transversal : la transformation se conduit en mode projet — comités de pilotage, référents métiers des deux filiales, prestataires externes, et le RSSI en sécurité transverse.
• Le « qui-fait-quoi » : l'axe vertical donne l'autorité, l'axe horizontal organise la coopération autour des projets (la RACI qui suit précise les rôles).
• Point d'attention : le cumul DSI/RSSI doit être surveillé (séparation des responsabilités sécurité) ; les ressources internes seront complétées par des prestataires.
""")
    return s


def _raci_cell(table, i, j, letter):
    fill = {"A": d.ORANGE, "R": TINT, "C": d.BLANC, "I": d.BLANC}[letter]
    color = {"A": d.BLANC, "R": d.ANTHRA, "C": GREY, "I": GREY}[letter]
    _cell(table.cell(i, j), letter, size=10, color=color, bold=letter in ("A", "R"),
          fill=fill, align=PP_ALIGN.CENTER)


def build_d10(prs):
    """D10 — Matrice RACI des responsabilités (parties prenantes)."""
    s = d.blank_slide(prs)
    d.header(s, "Matrice RACI des responsabilités",
             "Qui fait quoi sur la transformation du SI")

    cols = ["Activité / responsabilité", "DG", "DSI", "Chef de projet", "Équipes IT",
            "Métiers filiales", "Prestataires"]
    rows = [
        ("Stratégie & schéma directeur SI", "A", "R", "C", "I", "C", "I"),
        ("Budget & arbitrage des priorités", "A", "R", "C", "I", "C", "I"),
        ("Pilotage des projets de transformation", "I", "A", "R", "C", "C", "C"),
        ("Urbanisation & ERP unique (P7-P8)", "I", "A", "R", "C", "C", "C"),
        ("Gouvernance ITIL (services & incidents)", "I", "A", "C", "R", "I", "C"),
        ("Sécurité & conformité (SMSI / RGPD)", "I", "A", "I", "R", "I", "C"),
        ("Conduite du changement & formation", "I", "A", "R", "C", "C", "C"),
    ]

    X, Y, W = Inches(0.5), Inches(1.45), Inches(9.0)
    RH = Inches(0.34)
    n = len(rows) + 1
    gfx = s.shapes.add_table(n, len(cols), X, Y, W, RH * n)
    table = gfx.table
    table.first_row = False
    table.horz_banding = False
    table.columns[0].width = Inches(3.0)
    for j in range(1, len(cols)):
        table.columns[j].width = Inches(1.0)
    for r in table.rows:
        r.height = RH

    _cell(table.cell(0, 0), cols[0], size=10, color=d.BLANC, bold=True, fill=d.ANTHRA)
    for j in range(1, len(cols)):
        _cell(table.cell(0, j), cols[j], size=8.5, color=d.BLANC, bold=True,
              fill=d.ORANGE, align=PP_ALIGN.CENTER)
    for i, row in enumerate(rows, start=1):
        _cell(table.cell(i, 0), row[0], size=8.8, color=d.ANTHRA, bold=True, fill=d.CREME)
        for j in range(1, len(cols)):
            _raci_cell(table, i, j, row[j])

    tb, tf = d.textbox(s, X, Y + RH * n + Inches(0.12), W, Inches(0.35))
    d.set_para(tf.paragraphs[0],
               "R = Réalise   ·   A = Approuve (responsable, 1 seul par ligne)   ·   C = Consulté   ·   I = Informé",
               size=9.5, color=d.ORANGE, bold=True, align=PP_ALIGN.CENTER)

    d.footer(s, "RACI")
    d.notes(s, """
Matrice RACI — clarifie qui décide, qui fait, qui est consulté ou informé (critère fortement pondéré).
• Une seule lettre A (Approuve/responsable) par ligne : pas d'ambiguïté de décision.
• La DG approuve la stratégie, le budget et les arbitrages ; la DSI en est le réalisateur puis devient l'« A » sur l'exécution (projets, gouvernance, sécurité, changement).
• Les chefs de projet réalisent le pilotage et le changement ; les équipes IT réalisent la gouvernance ITIL et la sécurité.
• Les métiers des filiales sont consultés (expression de besoin, recette) et les prestataires consultés/contributeurs selon les lots.
• Insister : la RACI prévient les conflits de responsabilité, prérequis d'un pilotage multi-sites efficace.
""")
    return s


def _toolcards(s, items, y, ch):
    """Rangée de cartes outils (chip titre + description)."""
    n = len(items)
    gap = Inches(0.18)
    total = Inches(9.0)
    cw = (total - gap * (n - 1)) / n
    for i, (name, desc) in enumerate(items):
        x = Inches(0.5) + (cw + gap) * i
        d.card(s, x, y, cw, ch, fill=d.CREME)
        d.chip(s, x + Inches(0.1), y + Inches(0.12), cw - Inches(0.2), name, color=d.ORANGE, h=Inches(0.42))
        tb, tf = d.textbox(s, x + Inches(0.14), y + Inches(0.66), cw - Inches(0.28), ch - Inches(0.8))
        d.set_para(tf.paragraphs[0], desc, size=9.5, color=d.ANTHRA)


def build_d11(prs):
    """D11 — Outils de suivi des projets de transformation."""
    s = d.blank_slide(prs)
    d.header(s, "Outils de suivi des projets",
             "Piloter le portefeuille de transformation du SI")
    _toolcards(s, [
        ("Jira", "Suivi agile des projets de dév (CRM, e-commerce, BI) : backlog, sprints, tickets."),
        ("MS Project", "Planification des projets structurants (ERP P7-P8) : jalons, dépendances, charge."),
        ("Planner / Teams", "Tâches collaboratives et coordination des équipes au quotidien."),
        ("Power BI", "Tableau de bord du portefeuille : avancement, budget, KPI par projet."),
    ], Inches(1.6), Inches(2.4))

    tb, tf = d.textbox(s, Inches(0.5), Inches(4.35), Inches(9.0), Inches(0.7), anchor=MSO_ANCHOR.TOP)
    d.set_para(tf.paragraphs[0], "Pourquoi cet outillage", size=11, color=d.ORANGE, bold=True)
    d.add_para(tf, "Outils déjà présents dans l'écosystème Microsoft de Torpier (coût et adoption maîtrisés). Jira pour l'agile, MS Project pour le prédictif : on couvre les deux natures de projets, le tout consolidé dans Power BI pour la direction.",
               size=10, color=d.ANTHRA, space_before=3)
    d.footer(s, "Suivi projets")
    d.notes(s, """
Outils de suivi des projets — cohérents avec les référentiels et la nature des projets.
• Jira (agile/Scrum) pour les projets de développement, MS Project (prédictif) pour les projets structurants type ERP.
• Planner/Teams pour la coordination quotidienne, Power BI pour consolider l'avancement, le budget et les KPI à destination de la direction.
• Argument : on capitalise sur l'écosystème Microsoft déjà en place (maîtrise des coûts et de l'adoption) et on couvre les deux natures de projets.
""")
    return s


def build_d12(prs):
    """D12 — Outils de pilotage du changement."""
    s = d.blank_slide(prs)
    d.header(s, "Outils de pilotage du changement",
             "Mesurer l'adhésion et accompagner les utilisateurs")
    _toolcards(s, [
        ("Teams / Viva Engage", "Communication descendante et ascendante, animation de communauté."),
        ("MS Forms", "Baromètre d'adhésion : enquêtes avant/après, remontée des irritants."),
        ("LMS / e-learning", "Parcours de formation et montée en compétences des utilisateurs."),
        ("Power BI (adoption)", "Indicateurs d'adoption et d'usage des nouveaux outils."),
    ], Inches(1.6), Inches(2.4))

    tb, tf = d.textbox(s, Inches(0.5), Inches(4.35), Inches(9.0), Inches(0.7))
    d.set_para(tf.paragraphs[0], "Relais humains", size=11, color=d.ORANGE, bold=True)
    d.add_para(tf, "Au-delà des outils : un réseau d'ambassadeurs (un référent par site/filiale) et des points de communication réguliers. La mesure d'adoption nourrit la boucle d'amélioration (ADKAR / Kotter).",
               size=10, color=d.ANTHRA, space_before=3)
    d.footer(s, "Pilotage changement")
    d.notes(s, """
Outils de pilotage du changement — rendre l'adhésion mesurable.
• Teams/Viva Engage pour communiquer dans les deux sens ; MS Forms pour un baromètre d'adhésion (avant/après) ; LMS pour la formation ; Power BI pour suivre l'adoption.
• Relais humains : un réseau d'ambassadeurs par site/filiale, des points réguliers.
• Lien méthode : ces indicateurs alimentent la démarche de conduite du changement (ADKAR/Kotter) — on pilote l'humain comme un projet.
""")
    return s


def build_d13(prs):
    """D13 — Pitch d'adhésion à la direction."""
    s = d.blank_slide(prs)
    d.header(s, "Notre conviction pour la direction",
             "Pourquoi engager dès maintenant la conduite du changement")

    cols = [
        ("Le constat", d.ANTHRA, [
            "SI hétérogène, en silos entre filiales",
            "Support sans méthode ni indicateurs",
            "Infra fragile, sécurité non maîtrisée",
        ]),
        ("Notre proposition", d.ORANGE, [
            "Urbaniser le SI & ERP unique",
            "Cloud + gouvernance ITIL",
            "Sécuriser et mettre en conformité",
        ]),
        ("Les bénéfices", d.VERT, [
            "Productivité & données fiables",
            "Résilience et continuité de service",
            "SI aligné sur la stratégie 2026-2029",
        ]),
    ]
    CW, gap = Inches(2.92), Inches(0.12)
    Y, CH = Inches(1.4), Inches(2.7)
    for i, (titre, col, pts) in enumerate(cols):
        x = Inches(0.5) + (CW + gap) * i
        d.card(s, x, Y, CW, CH, fill=d.CREME)
        d.chip(s, x + Inches(0.12), Y + Inches(0.12), CW - Inches(0.24), titre, color=col, h=Inches(0.42))
        tb, tf = d.textbox(s, x + Inches(0.18), Y + Inches(0.68), CW - Inches(0.36), CH - Inches(0.8))
        d.set_para(tf.paragraphs[0], "•  " + pts[0], size=10.5, color=d.ANTHRA)
        for p in pts[1:]:
            d.add_para(tf, "•  " + p, size=10.5, color=d.ANTHRA, space_before=6)

    ask = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.35), Inches(9.0), Inches(0.95))
    ask.fill.solid(); ask.fill.fore_color.rgb = d.ORANGE
    ask.line.fill.background(); ask.shadow.inherit = False
    ask.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    ask.text_frame.word_wrap = True
    d.set_para(ask.text_frame.paragraphs[0],
               "Notre demande : valider la politique de conduite du changement et lancer le schéma directeur SI.",
               size=14, color=d.BLANC, bold=True, align=PP_ALIGN.CENTER)
    d.add_para(ask.text_frame, "Sans accompagnement du changement, même le meilleur projet SI échoue à l'usage.",
               size=10, color=d.BLANC, align=PP_ALIGN.CENTER, space_before=3)

    d.footer(s, "Pitch")
    d.notes(s, """
Pitch d'adhésion — le moment de convaincre la direction (critère fortement pondéré).
• Dérouler la logique constat → proposition → bénéfices, en restant sur le « pourquoi ».
• Constat : les faiblesses déjà démontrées (silos, support informel, infra/sécurité fragiles).
• Proposition : nos axes (urbanisation/ERP, cloud + ITIL, sécurité/conformité).
• Bénéfices : productivité, fiabilité des données, résilience, alignement stratégique — le langage de la direction (valeur et risque).
• L'ASK : faire valider la politique de conduite du changement et lancer le schéma directeur. Message clé : la technique ne suffit pas, l'adhésion des équipes conditionne la réussite.
""")
    return s


def build_d14(prs):
    """D14 — Organisation du groupe de travail (rôles & responsabilités)."""
    s = d.blank_slide(prs)
    d.header(s, "Organisation du groupe de travail",
             "Rôle et contribution de chaque membre à l'étude")

    data = [
        ("[Membre 1]", "Pilote / coordination", "Cadrage, planning, cohérence collectif ↔ individuel"),
        ("[Membre 2]", "Analyse de l'existant", "Cartographie processus, tableaux applis & infra"),
        ("[Membre 3]", "Gouvernance & équipes", "Référentiels, organisation SI, matrice RACI"),
        ("[Membre 4]", "Changement & outils", "Outils de suivi/changement, pitch d'adhésion"),
    ]
    X, Y, W = Inches(0.5), Inches(1.6), Inches(9.0)
    RH = Inches(0.5)
    n = len(data) + 1
    gfx = s.shapes.add_table(n, 3, X, Y, W, RH * n)
    table = gfx.table
    table.first_row = False
    table.horz_banding = False
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(2.8)
    table.columns[2].width = Inches(4.0)
    for r in table.rows:
        r.height = RH
    for j, h in enumerate(("Membre", "Rôle dans le groupe", "Contribution principale")):
        _cell(table.cell(0, j), h, size=11, color=d.BLANC, bold=True,
              fill=d.ANTHRA if j == 0 else d.ORANGE)
    for i, (m, role, contrib) in enumerate(data, start=1):
        _cell(table.cell(i, 0), m, size=10, color=d.ANTHRA, bold=True, fill=d.CREME)
        _cell(table.cell(i, 1), role, size=10, color=d.ANTHRA, fill=d.BLANC)
        _cell(table.cell(i, 2), contrib, size=10, color=d.ANTHRA, fill=d.BLANC)

    tb, tf = d.textbox(s, X, Y + RH * n + Inches(0.15), W, Inches(0.4))
    d.set_para(tf.paragraphs[0],
               "Travail co-construit « au fil de l'eau » : production partagée, revues croisées, jeu de données commun pour garantir la cohérence.",
               size=9.5, color=d.ORANGE, bold=True)
    d.footer(s, "Groupe")
    d.notes(s, """
Organisation du groupe — montrer une équipe structurée (à compléter avec les vrais noms/rôles).
• Chaque membre porte un volet (pilotage, analyse de l'existant, gouvernance, changement) mais le livrable est co-construit.
• Insister sur la méthode : revues croisées et jeu de données commun → cohérence entre le collectif et les soutenances individuelles.
• ⚠️ Remplacer [Membre 1-4] par les noms réels et ajuster les rôles selon la répartition effective.
""")
    return s


def build_d15(prs):
    """D15 — Bilan du groupe (points forts / axes d'amélioration)."""
    s = d.blank_slide(prs)
    d.header(s, "Bilan du groupe", "Ce qui a bien fonctionné et nos axes de progrès")

    LX, RX, W = Inches(0.5), Inches(5.15), Inches(4.35)
    Y, CH = Inches(1.5), Inches(3.4)
    d.card(s, LX, Y, W, CH, fill=d.CREME)
    d.card(s, RX, Y, W, CH, fill=d.CREME)
    d.chip(s, LX + Inches(0.18), Y + Inches(0.18), Inches(2.6), "Points forts", color=d.VERT)
    d.chip(s, RX + Inches(0.18), Y + Inches(0.18), Inches(2.6), "Axes d'amélioration", color=d.ORANGE)

    tb, tf = d.textbox(s, LX + Inches(0.2), Y + Inches(0.66), W - Inches(0.4), CH - Inches(0.8))
    for k, t in enumerate([
        "Données réelles exploitées (SDSI, témoignages) → propositions crédibles",
        "Cohérence forte entre le collectif et les livrables individuels",
        "Argumentation systématique adossée aux référentiels (ITIL/COBIT)",
        "Répartition claire des rôles et production au fil de l'eau",
    ]):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        d.set_para(p, "•  " + t, size=10.5, color=d.ANTHRA, space_before=0 if k == 0 else 8)

    tb2, tf2 = d.textbox(s, RX + Inches(0.2), Y + Inches(0.66), W - Inches(0.4), CH - Inches(0.8))
    for k, t in enumerate([
        "Affiner le chiffrage et le ROI des projets",
        "Mieux cadrer le temps consacré à chaque livrable",
        "Approfondir les liens entre les filiales (POTIER)",
        "Renforcer la gestion des risques projet",
    ]):
        p = tf2.paragraphs[0] if k == 0 else tf2.add_paragraph()
        d.set_para(p, "•  " + t, size=10.5, color=d.ANTHRA, space_before=0 if k == 0 else 8)

    d.footer(s, "Bilan")
    d.notes(s, """
Bilan du groupe — exercice de recul honnête et constructif.
• Points forts : exploitation des données réelles (SDSI, témoignages), cohérence collectif/individuel, argumentation par les référentiels, organisation claire.
• Axes d'amélioration : chiffrage/ROI à affiner, meilleure gestion du temps, approfondir POTIER, renforcer la gestion des risques.
• Ton : montrer une équipe lucide et apprenante — c'est aussi une posture de manager.
""")
    return s


def build_d16(prs):
    """D16 — Conclusion + transition Q/R."""
    s = d.blank_slide(prs)
    d.header(s, "Conclusion — la trajectoire du SI Torpier",
             "Une transformation par étapes, alignée sur la stratégie")

    horizons = [
        ("Court terme", "ERP lot 1 (finances), gestion des incidents ITIL, audit de sécurité"),
        ("Moyen terme", "Move to Cloud, ERP lot 2 (achats/stocks), décisionnel groupe"),
        ("Long terme", "SI urbanisé, agile et sécurisé, pleinement aligné sur la stratégie"),
    ]
    CW, gap = Inches(2.92), Inches(0.12)
    Y, CH = Inches(1.5), Inches(2.2)
    for i, (t, desc) in enumerate(horizons):
        x = Inches(0.5) + (CW + gap) * i
        d.card(s, x, Y, CW, CH, fill=d.CREME)
        num_badge(s, x + CW / 2 - Inches(0.22), Y + Inches(0.18), i + 1)
        tb, tf = d.textbox(s, x + Inches(0.16), Y + Inches(0.78), CW - Inches(0.32), CH - Inches(0.9))
        d.set_para(tf.paragraphs[0], t, size=13, color=d.ORANGE, bold=True, align=PP_ALIGN.CENTER)
        d.add_para(tf, desc, size=10, color=d.ANTHRA, align=PP_ALIGN.CENTER, space_before=4)
        if i < 2:
            cb, cf = d.textbox(s, x + CW - Inches(0.02), Y + CH / 2 - Inches(0.25), gap + Inches(0.06),
                               Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
            d.set_para(cf.paragraphs[0], "›", size=22, color=d.ORANGE, bold=True, align=PP_ALIGN.CENTER)

    band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.2), Inches(9.0), Inches(0.95))
    band.fill.solid(); band.fill.fore_color.rgb = d.ANTHRA
    band.line.fill.background(); band.shadow.inherit = False
    band.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    d.set_para(band.text_frame.paragraphs[0],
               "Merci de votre attention — place à vos questions",
               size=16, color=d.BLANC, bold=True, align=PP_ALIGN.CENTER)
    d.add_para(band.text_frame, "Un SI piloté, urbanisé et sécurisé : le moteur de la stratégie Torpier 2026-2029.",
               size=10.5, color=d.GRIS, align=PP_ALIGN.CENTER, space_before=3)

    d.footer(s, "Conclusion")
    d.notes(s, """
Conclusion — donner une vision claire et ouvrir l'échange.
• Résumer la trajectoire en 3 horizons : quick wins (ERP lot 1, incidents ITIL, audit sécurité), moyen terme (cloud, ERP lot 2, BI groupe), cible (SI urbanisé, agile, sécurisé).
• Reboucler sur la stratégie 2026-2029 : le SI est le moteur de la croissance, du durable et de l'industrie 4.0.
• Inviter aux questions avec confiance ; garder en réserve les notes des slides précédentes pour les Q/R.
""")
    return s


def main():
    prs = Presentation(SRC)
    title = prs.slides[TPL["title"]]
    swot = prs.slides[TPL["swot"]]

    build_d1(title)
    d2 = build_d2(prs)
    build_d3(swot)
    d4 = build_d4(prs)
    d5 = build_d5(prs)
    d6 = build_d6(prs)
    d7 = build_d7(prs)
    d8 = build_d8(prs)
    d9 = build_d9(prs)
    d10 = build_d10(prs)
    d11 = build_d11(prs)
    d12 = build_d12(prs)
    d13 = build_d13(prs)
    d14 = build_d14(prs)
    d15 = build_d15(prs)
    d16 = build_d16(prs)

    reorder_and_prune(prs, [title, d2, swot, d4, d5, d6, d7, d8, d9, d10,
                            d11, d12, d13, d14, d15, d16])
    prs.save(OUT)
    print("Slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
